from pptx import Presentation
import sys

# Load the existing template
template_path = sys.argv[1]
prs = Presentation(template_path)

# Delete existing slides
xml_slides = prs.slides._sldIdLst
slides = list(xml_slides)
for slide in slides:
    xml_slides.remove(slide)

def add_slide(prs, layout_idx, title_text, content_text=None):
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title_text
        
    # Set content
    if content_text is not None and len(slide.placeholders) > 1:
        # Find the first placeholder that isn't the title
        body_shape = None
        for shape in slide.placeholders:
            if shape != slide.shapes.title and shape.has_text_frame:
                body_shape = shape
                break
                
        if body_shape:
            tf = body_shape.text_frame
            if isinstance(content_text, list):
                for i, point in enumerate(content_text):
                    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                    p.text = point
                    p.level = 0
                    if p.text.startswith("  -"):
                        p.text = p.text[3:]
                        p.level = 1
            else:
                tf.text = content_text
    return slide

# Slide 1: Title Slide (Layout 0)
slide = add_slide(prs, 0, "NPKMによる再現性と品質を担保したインフラ自動化環境の構築")
if len(slide.placeholders) > 1:
    subtitle = slide.placeholders[1]
    if subtitle.has_text_frame:
        subtitle.text = "Nuke Playbook Kit Manager\nKarabiner Software LLC"

# Slide 2: Objective (Layout 1)
add_slide(prs, 1, "本資料の目的", 
    ["インフラ構築・アプリケーションデプロイの再現性・正確性を担保するため、",
     "Playbook定義・変数・実行環境をネイティブかつ依存ゼロな状態で完全に管理すること"])

# Slide 3: Need for Management
add_slide(prs, 1, "インフラ構成管理の必要性",
    ["システムは複数の要素から構成される",
     "手動作業やシェルスクリプトだけでは、以下の構成要素の完全な再現が困難：",
     "  - サーバーのOS設定やミドルウェアのインストール",
     "  - 環境ごとの変数（Dev/Stg/Prod）",
     "  - デプロイ手順やオーケストレーション",
     "これらを「コード（Infrastructure as Code）」として管理する必要がある"])

# Slide 4: Problems
add_slide(prs, 1, "自動化が不十分な場合の問題",
    ["手順書による手動構築や、複雑化したスクリプトに依存すると、",
     "  - 同じ環境を二度と再現できない（冪等性の欠如）",
     "  - 構築に時間がかかりすぎる",
     "  - 障害時に「いつ・誰が・何を変えたか」追えない",
     "  - Python等、実行基盤自体のバージョン依存で自動化ツール自体が動かない",
     "という深刻な問題が発生する"])

# Slide 5: Ideal
add_slide(prs, 1, "目指す姿（理想）",
    ["Release Versionに対して、インフラ構築手順（Playbook）と変数が紐づいており、",
     "どの環境に対しても、高速かつ確実に、完全に同一のシステム構成を再現可能な状態",
     "",
     "ソースコード ＋ インフラ構成 ＋ NPKM ＝ 100%の再現性"])

# Slide 6: Issues with Traditional Tools
add_slide(prs, 1, "既存自動化ツール（Ansible等）の課題",
    ["Ansible等は普及しているが、以下の技術的制約がある：",
     "  - Python依存：実行元ノードに適切なPython環境とライブラリが必須",
     "  - 実行速度：マルチプロセスのForkオーバーヘッドが大きく、並列処理が遅い",
     "  - 配布の難しさ：自動化ツール自体をセットアップする手間がかかる",
     "品質を担保した検証環境を「短期間で構築」する際のボトルネックになり得る"])

# Slide 7: Solution - NPKM
add_slide(prs, 1, "解決策：NPKMの採用",
    ["既存のAnsible Playbook資産をそのまま活用しながら、",
     "実行基盤を「完全ネイティブ化」するNPKM（Nuke Playbook Kit Manager）を採用",
     "",
     "NPKMは、Coni言語（Goの強力なDSLとして動作）で実装されたAnsible互換エンジンであり、",
     "依存ゼロの単一バイナリとして、既存の構成管理の課題を全て解決する"])

# Slide 8: Technical Reasons
add_slide(prs, 1, "NPKMを採用する技術的な理由",
    ["① 依存関係ゼロ",
     "  - Python環境不要。単一の実行可能バイナリ（Mac/Linux/Win）を置くだけで動作",
     "② 圧倒的な高速性と並列処理",
     "  - Goroutine（Go言語の軽量スレッド）による高速なFan-out / Fan-in並列処理",
     "③ Coni言語による柔軟性",
     "  - GoのDSLとして構築されているため、Goの強力なエコシステムと直接統合可能",
     "④ 高度なセキュリティ",
     "  - Ansible Vault互換の暗号化モジュールをバイナリに標準内蔵"])

# Slide 9: Best Features
add_slide(prs, 1, "NPKMの強力なベスト機能群",
    ["NPKMは単なる実行エンジンにとどまらない、強力な開発者支援機能を備える：",
     "  - 自動ドキュメント生成 (--doc): Playbookの実行フローをMermaid図として出力し、処理の可視化を実現",
     "  - インタラクティブ・デバッグ (--step): タスクごとに変数の展開結果を確認しながらステップ実行が可能",
     "  - 静的解析エンジン内蔵 (npkm-lint): 実行前にPlaybookの構文・引数エラーを即座に検知",
     "  - 並列タスクブロック: ホスト内のタスクすらも並列実行可能な parallel ディレクティブ"])

# Slide 10: Positioning
add_slide(prs, 1, "NPKMの位置づけ",
    ["NPKMは「自動化の手法」をゼロから作り直すものではない。",
     "",
     "業界標準（AnsibleベースのPlaybook）の記述モデルを維持したまま、",
     "「実行エンジンのアーキテクチャ」だけを次世代（高速・安全・依存ゼロ）へとアップグレードする、",
     "堅牢なインフラストラクチャ・オーケストレーション基盤である。"])

# Slide 11: Summary
add_slide(prs, 1, "まとめ（導入方針）",
    ["① 既存のAnsible Playbook資産を再利用し、NPKMへシームレスに移行",
     "② Python環境構築の手間を排除し、単一バイナリでどこでも即座に自動化を開始",
     "③ 高速な並列実行により、ユーザー検証環境や本番環境の構築リードタイムを大幅に短縮",
     "④ 実行結果の再現性と正当性を確実なものにし、インフラの品質を極限まで高める"])

out_path = '/Users/nico/Downloads/NPKM_Reproducibility_and_Quality_Assurance.pptx'
prs.save(out_path)
print(f"Successfully generated {out_path}")
