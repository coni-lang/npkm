from pptx import Presentation
import sys

def extract_text(filename):
    prs = Presentation(filename)
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                print(paragraph.text)
        print()

extract_text(sys.argv[1])
